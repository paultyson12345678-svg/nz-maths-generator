# exporters.py
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from PIL import Image, ImageDraw
import io

def generate_powerpoint_slide(task_title, task_body, questions, extension, phase, theme):
    """Generates a 16:9 PowerPoint presentation slide downloadable as .pptx."""
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)
    
    blank_slide_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(blank_slide_layout)
    
    # Title Box
    txBox = slide.shapes.add_textbox(Inches(0.8), Inches(0.6), Inches(11.7), Inches(1.0))
    tf = txBox.text_frame
    p = tf.paragraphs[0]
    p.text = task_title
    p.font.size = Pt(28)
    p.font.bold = True
    p.font.color.rgb = RGBColor(16, 44, 87) # Deep Blue
    
    # Subtitle / Context Tag
    p2 = tf.add_paragraph()
    p2.text = f"Context: {theme} | Alignment: {phase}"
    p2.font.size = Pt(14)
    p2.font.color.rgb = RGBColor(100, 100, 100)
    
    # Task Body Box
    body_box = slide.shapes.add_textbox(Inches(0.8), Inches(1.8), Inches(11.7), Inches(2.0))
    bf = body_box.text_frame
    bf.word_wrap = True
    p_body = bf.paragraphs[0]
    p_body.text = task_body
    p_body.font.size = Pt(18)
    
    # Questions Box
    q_box = slide.shapes.add_textbox(Inches(0.8), Inches(4.0), Inches(11.7), Inches(2.8))
    qf = q_box.text_frame
    qf.word_wrap = True
    
    p_q_head = qf.paragraphs[0]
    p_q_head.text = "Task Questions:"
    p_q_head.font.bold = True
    p_q_head.font.size = Pt(18)
    
    for idx, q in enumerate(questions, 1):
        pq = qf.add_paragraph()
        pq.text = f"{idx}. {q}"
        pq.font.size = Pt(16)
        
    # Extension Challenge
    if extension:
        pe = qf.add_paragraph()
        pe.text = f"\n🚀 Extension Challenge: {extension}"
        pe.font.size = Pt(16)
        pe.font.bold = True
        pe.font.color.rgb = RGBColor(0, 128, 96) # Green
        
    output = io.BytesIO()
    prs.save(output)
    output.seek(0)
    return output

def generate_task_card_image(task_title, task_body, questions, extension):
    """Creates a downloadable PNG task card."""
    img = Image.new('RGB', (1000, 700), color=(250, 252, 255))
    draw = ImageDraw.Draw(img)
    
    # Outer Border & Header Bar
    draw.rectangle([20, 20, 980, 680], outline=(40, 80, 140), width=4)
    draw.rectangle([20, 20, 980, 90], fill=(40, 80, 140))
    
    # Title Header
    draw.text((40, 40), task_title, fill=(255, 255, 255))
    
    y_offset = 120
    draw.text((40, y_offset), "Scenario:", fill=(0, 0, 0))
    y_offset += 30
    
    # Basic word wrapping for task body
    words = task_body.split()
    line = ""
    for word in words:
        if len(line + " " + word) < 70:
            line += " " + word
        else:
            draw.text((40, y_offset), line, fill=(40, 40, 40))
            y_offset += 25
            line = word
    draw.text((40, y_offset), line, fill=(40, 40, 40))
    y_offset += 45
    
    # Questions
    draw.text((40, y_offset), "Questions:", fill=(0, 0, 0))
    y_offset += 30
    for idx, q in enumerate(questions, 1):
        draw.text((50, y_offset), f"{idx}. {q}", fill=(30, 30, 30))
        y_offset += 35
        
    # Extension Question
    if extension:
        y_offset += 15
        draw.text((40, y_offset), f"Extension: {extension}", fill=(0, 120, 60))
        
    img_byte_arr = io.BytesIO()
    img.save(img_byte_arr, format='PNG')
    img_byte_arr.seek(0)
    return img_byte_arr