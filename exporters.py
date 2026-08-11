import io
import os
import re
from reportlab.lib.pagesizes import A4
from reportlab.platypus import SimpleDocTemplate, Paragraph, Spacer, Table, TableStyle, PageBreak
from reportlab.lib.styles import getSampleStyleSheet, ParagraphStyle
from reportlab.lib import colors
from reportlab.pdfgen import canvas
from reportlab.pdfbase import pdfmetrics
from reportlab.pdfbase.ttfonts import TTFont

# Optional pptx import for slide export
try:
    from pptx import Presentation
    from pptx.util import Inches, Pt
    from pptx.dml.color import RGBColor
    from pptx.enum.text import MSO_ANCHOR
    PPTX_AVAILABLE = True
except ImportError:
    PPTX_AVAILABLE = False


# --- FONT REGISTRATION ---
_font_registered = False
_registered_font_name = 'Helvetica'
_registered_bold_font_name = 'Helvetica-Bold'

def get_macron_font():
    global _font_registered, _registered_font_name, _registered_bold_font_name
    if _font_registered:
        return _registered_font_name, _registered_bold_font_name

    local_font_path = os.path.join(os.path.dirname(__file__), 'DejaVuSans.ttf')

    font_paths = [
        local_font_path,
        '/usr/share/fonts/truetype/dejavu/DejaVuSans.ttf',
        '/usr/share/fonts/TTF/DejaVuSans.ttf',
        '/Library/Fonts/DejaVuSans.ttf'
    ]

    for path in font_paths:
        if os.path.exists(path):
            try:
                pdfmetrics.registerFont(TTFont('DejaVuSans', path))
                _registered_font_name = 'DejaVuSans'
                _registered_bold_font_name = 'DejaVuSans'
                _font_registered = True
                break
            except Exception:
                pass

    return _registered_font_name, _registered_bold_font_name


# --- DATA CLEANUP HELPER ---
def _clean_and_extract_data(answers, teacher_notes, misconceptions):
    """Intercepts and separates lumped text if the AI failed to format it properly."""
    extracted_tn = teacher_notes
    extracted_mc = misconceptions
    
    if not answers:
        return answers, extracted_tn, extracted_mc

    def process_string(text):
        nonlocal extracted_tn, extracted_mc
        if not text: return text
        
        # Look for Teacher Guidance stuck in the answer
        if "Teacher Guidance:" in text:
            parts = text.split("Teacher Guidance:")
            text = parts[0].strip()
            
            # Check if Misconceptions is also bundled in there
            if "Common Misconceptions:" in parts[1]:
                sub_parts = parts[1].split("Common Misconceptions:")
                extracted_tn = sub_parts[0].strip() if not extracted_tn else extracted_tn
                extracted_mc = sub_parts[1].strip() if not extracted_mc else extracted_mc
            else:
                extracted_tn = parts[1].strip() if not extracted_tn else extracted_tn
        
        # Look for Misconceptions stuck in the answer
        if "Common Misconceptions:" in text:
            parts = text.split("Common Misconceptions:")
            text = parts[0].strip()
            
            if "Teacher Guidance:" in parts[1]:
                sub_parts = parts[1].split("Teacher Guidance:")
                extracted_mc = sub_parts[0].strip() if not extracted_mc else extracted_mc
                extracted_tn = sub_parts[1].strip() if not extracted_tn else extracted_tn
            else:
                extracted_mc = parts[1].strip() if not extracted_mc else extracted_mc
                
        return text

    if isinstance(answers, list):
        cleaned_answers = [process_string(str(ans)) for ans in answers]
    else:
        cleaned_answers = process_string(str(answers))
        
    return cleaned_answers, extracted_tn, extracted_mc


# --- POWERPOINT GENERATOR ---
def generate_powerpoint_slide(title, scenario, questions, extension, phase, theme, answers=None, teacher_notes=None, misconceptions=None):
    """
    Generates a PowerPoint presentation (.pptx) with the math task.
    """
    if not PPTX_AVAILABLE:
        return None
        
    # --- Intercept and clean the data before making the slides ---
    answers, teacher_notes, misconceptions = _clean_and_extract_data(answers, teacher_notes, misconceptions)

    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)
    blank_layout = prs.slide_layouts[6]

    # --- SLIDE 1: Scenario & Question 1 ---
    slide1 = prs.slides.add_slide(blank_layout)

    title_box1 = slide1.shapes.add_textbox(Inches(0.8), Inches(0.4), Inches(11.7), Inches(0.7))
    tf1 = title_box1.text_frame
    tf1.word_wrap = True
    p1 = tf1.paragraphs[0]
    p1.text = title
    p1.font.size = Pt(24)
    p1.font.bold = True
    p1.font.color.rgb = RGBColor(30, 58, 138)

    scen_box = slide1.shapes.add_textbox(Inches(0.8), Inches(1.2), Inches(11.7), Inches(2.0))
    tf_scen = scen_box.text_frame
    tf_scen.word_wrap = True
    p_scen = tf_scen.paragraphs[0]
    p_scen.text = f"Scenario: {scenario}"
    p_scen.font.size = Pt(18)
    p_scen.font.color.rgb = RGBColor(31, 41, 55)

    q1_box = slide1.shapes.add_textbox(Inches(0.8), Inches(3.6), Inches(11.7), Inches(3.2))
    tf_q1 = q1_box.text_frame
    tf_q1.word_wrap = True
    if questions:
        p_q1 = tf_q1.paragraphs[0]
        p_q1.text = f"1. {questions[0]}"
        p_q1.font.size = Pt(20)
        p_q1.font.color.rgb = RGBColor(30, 41, 59)

    # --- SLIDE 2: Question 2+ & Extension Challenge ---
    slide2 = prs.slides.add_slide(blank_layout)

    title_box2 = slide2.shapes.add_textbox(Inches(0.8), Inches(0.4), Inches(11.7), Inches(0.7))
    tf2 = title_box2.text_frame
    tf2.word_wrap = True
    p2 = tf2.paragraphs[0]
    p2.text = f"{title} (Continued)"
    p2.font.size = Pt(24)
    p2.font.bold = True
    p2.font.color.rgb = RGBColor(30, 58, 138)

    q2_box = slide2.shapes.add_textbox(Inches(0.8), Inches(1.4), Inches(11.7), Inches(5.5))
    tf_q2 = q2_box.text_frame
    tf_q2.word_wrap = True

    remaining_qs = questions[1:] if len(questions) > 1 else []
    first_item = True

    for idx, q in enumerate(remaining_qs, 2):
        p_q = tf_q2.paragraphs[0] if first_item else tf_q2.add_paragraph()
        first_item = False
        p_q.text = f"{idx}. {q}"
        p_q.font.size = Pt(20)
        p_q.space_after = Pt(30) 

    if extension:
        p_ext = tf_q2.paragraphs[0] if first_item else tf_q2.add_paragraph()
        p_ext.text = f"Extension Challenge:\n{extension}"
        p_ext.font.size = Pt(20)
        p_ext.font.bold = True
        p_ext.font.color.rgb = RGBColor(180, 83, 9)
        p_ext.space_before = Pt(40)

    # --- SLIDE 3: TEACHER NOTES, MISCONCEPTIONS & SOLUTIONS ---
    if answers or misconceptions or teacher_notes:
        slide_answers = prs.slides.add_slide(blank_layout)
        
        title_box3 = slide_answers.shapes.add_textbox(Inches(0.8), Inches(0.4), Inches(11.7), Inches(0.7))
        tf3 = title_box3.text_frame
        tf3.word_wrap = True
        title_p = tf3.paragraphs[0]
        title_p.text = f"Solutions & Notes: {title}"
        title_p.font.size = Pt(24)
        title_p.font.bold = True
        title_p.font.color.rgb = RGBColor(30, 58, 138)
        
        content_box = slide_answers.shapes.add_textbox(Inches(0.8), Inches(1.3), Inches(11.7), Inches(5.8))
        tf_content = content_box.text_frame
        tf_content.word_wrap = True
        tf_content.vertical_anchor = MSO_ANCHOR.TOP 
        
        is_first = True

        # 1. Teacher Guidance FIRST
        if teacher_notes:
            p_tn_title = tf_content.paragraphs[0] if is_first else tf_content.add_paragraph()
            is_first = False
            p_tn_title.text = "Teacher Guidance:"
            p_tn_title.font.bold = True
            p_tn_title.font.size = Pt(14)
            
            p_tn = tf_content.add_paragraph()
            p_tn.text = str(teacher_notes)
            p_tn.font.size = Pt(14)
            p_tn.space_after = Pt(18)
        
        # 2. Common Misconceptions SECOND
        if misconceptions:
            p_mc_title = tf_content.paragraphs[0] if is_first else tf_content.add_paragraph()
            is_first = False
            p_mc_title.text = "Common Misconceptions:"
            p_mc_title.font.bold = True
            p_mc_title.font.size = Pt(14)
            
            p_mc = tf_content.add_paragraph()
            p_mc.text = str(misconceptions)
            p_mc.font.size = Pt(14)
            p_mc.space_after = Pt(18)
            
        # 3. Solutions & Extension THIRD
        if answers:
            if isinstance(answers, list):
                for i, ans in enumerate(answers):
                    # Skip empty answers that might have resulted from cleanup
                    if not str(ans).strip(): 
                        continue
                        
                    p_head = tf_content.paragraphs[0] if is_first else tf_content.add_paragraph()
                    is_first = False
                    
                    if i == len(answers) - 1:
                        p_head.text = "Extension Solution:"
                    else:
                        p_head.text = f"Question {i+1}:"
                    
                    p_head.font.bold = True
                    p_head.font.size = Pt(14)
                    
                    p_ans = tf_content.add_paragraph()
                    p_ans.text = str(ans)
                    p_ans.font.size = Pt(14)
                    p_ans.space_after = Pt(12)
            else:
                if str(answers).strip():
                    p_head = tf_content.paragraphs[0] if is_first else tf_content.add_paragraph()
                    is_first = False
                    p_head.text = "Solutions:"
                    p_head.font.bold = True
                    p_head.font.size = Pt(14)
                    
                    p_ans = tf_content.add_paragraph()
                    p_ans.text = str(answers)
                    p_ans.font.size = Pt(14)
                    p_ans.space_after = Pt(12)

    ppt_buffer = io.BytesIO()
    prs.save(ppt_buffer)
    ppt_buffer.seek(0)
    return ppt_buffer.getvalue()


# --- PDF GENERATOR ---
def generate_task_pdf(title, scenario, questions, extension, phase, theme, answers=None, teacher_notes=None, misconceptions=None):
    """
    Generates a printable A4 PDF student worksheet with a solutions & notes page.
    """
    buffer = io.BytesIO()
    font_normal, font_bold = get_macron_font()
    
    # --- Intercept and clean the data before making the PDF ---
    answers, teacher_notes, misconceptions = _clean_and_extract_data(answers, teacher_notes, misconceptions)

    doc = SimpleDocTemplate(
        buffer,
        pagesize=A4,
        rightMargin=36,
        leftMargin=36,
        topMargin=24,
        bottomMargin=24
    )

    story = []
    styles = getSampleStyleSheet()

    total_text_length = len(title) + len(scenario) + sum(len(q) for q in questions) + (len(extension) if extension else 0)
    
    if total_text_length > 600:
        title_size, title_lead = 15, 19
        body_size, body_lead = 9, 12
        q_size, q_lead = 9.5, 13
        scen_padding = 5
        space_after_box = 10
    else:
        title_size, title_lead = 17, 21
        body_size, body_lead = 9.5, 13
        q_size, q_lead = 10, 14
        scen_padding = 7
        space_after_box = 12

    title_style = ParagraphStyle(
        'DocTitle',
        parent=styles['Heading1'],
        fontName=font_bold,
        fontSize=title_size,
        leading=title_lead,
        textColor=colors.HexColor('#1B365D'),
        spaceAfter=2
    )

    meta_style = ParagraphStyle(
        'DocMeta',
        parent=styles['Normal'],
        fontName=font_bold,
        fontSize=9,
        leading=11,
        textColor=colors.HexColor('#4A5568')
    )

    scenario_style = ParagraphStyle(
        'ScenarioText',
        parent=styles['Normal'],
        fontName=font_normal,
        fontSize=body_size,
        leading=body_lead,
        textColor=colors.HexColor('#2D3748')
    )

    question_style = ParagraphStyle(
        'QuestionText',
        parent=styles['Normal'],
        fontName=font_bold,
        fontSize=q_size,
        leading=q_lead,
        textColor=colors.HexColor('#1A202C')
    )

    story.append(Paragraph(title, title_style))
    story.append(Paragraph(f"<b>Phase:</b> {phase} &nbsp;&nbsp;|&nbsp;&nbsp; <b>Context:</b> {theme}", meta_style))
    story.append(Spacer(1, 4))

    scenario_p = Paragraph(f"<b>Scenario:</b><br/>{scenario}", scenario_style)
    scenario_table = Table([[scenario_p]], colWidths=[523])
    scenario_table.setStyle(TableStyle([
        ('BACKGROUND', (0, 0), (-1, -1), colors.HexColor('#F7FAFC')),
        ('BOX', (0, 0), (-1, -1), 1, colors.HexColor('#CBD5E0')),
        ('PADDING', (0, 0), (-1, -1), scen_padding),
        ('VALIGN', (0, 0), (-1, -1), 'TOP'),
    ]))
    story.append(scenario_table)
    story.append(Spacer(1, 10))

    num_questions = len(questions)
    total_units = num_questions + (2 if extension else 0)
    
    estimated_text_lines = sum(max(1, len(q) // 80) for q in questions) + (max(1, len(extension) // 80) if extension else 0)
    text_height = estimated_text_lines * q_lead
    
    gaps_space = (num_questions + (1 if extension else 0)) * (space_after_box + 4)
    available_box_space = 460 - text_height - gaps_space
    
    unit_height = max(35, min(75, available_box_space / max(1, total_units)))
    standard_box_height = unit_height
    extension_box_height = unit_height * 2 

    def create_working_box(box_height):
        t = Table([['']], colWidths=[523], rowHeights=[box_height])
        t.setStyle(TableStyle([
            ('BOX', (0, 0), (-1, -1), 1, colors.HexColor('#A0AEC0')),
            ('BACKGROUND', (0, 0), (-1, -1), colors.white),
        ]))
        return t

    for idx, q_text in enumerate(questions, 1):
        story.append(Paragraph(f"<b>Question {idx}:</b> {q_text}", question_style))
        story.append(Spacer(1, 3))
        story.append(create_working_box(box_height=standard_box_height))
        story.append(Spacer(1, space_after_box))

    if extension:
        story.append(Paragraph(f"<b>Extension Challenge:</b> {extension}", question_style))
        story.append(Spacer(1, 3))
        story.append(create_working_box(box_height=extension_box_height))
        story.append(Spacer(1, space_after_box))

    # --- PAGE 2: TEACHER NOTES, MISCONCEPTIONS & SOLUTIONS ---
    if answers or teacher_notes or misconceptions:
        story.append(PageBreak()) 
        
        section_title_style = ParagraphStyle(
            'SectionTitle',
            parent=styles['Heading1'],
            fontName=font_bold,
            fontSize=18,
            leading=22,
            textColor=colors.HexColor('#1B365D'),
            spaceAfter=15
        )
        
        body_style = ParagraphStyle(
            'SectionBody',
            parent=styles['Normal'],
            fontName=font_normal,
            fontSize=12,
            leading=16,
            textColor=colors.HexColor('#2D3748'),
            spaceAfter=12
        )
        
        story.append(Paragraph("Solutions & Notes", section_title_style))

        if teacher_notes:
            story.append(Paragraph("<b>Teacher Guidance:</b>", body_style))
            story.append(Paragraph(str(teacher_notes), body_style))
            story.append(Spacer(1, 15)) 
            
        if misconceptions:
            story.append(Paragraph("<b>Common Misconceptions:</b>", body_style))
            story.append(Paragraph(str(misconceptions), body_style))
            story.append(Spacer(1, 4))
            
        if answers:
            if isinstance(answers, list):
                for i, ans in enumerate(answers):
                    if not str(ans).strip():
                        continue
                    if i == len(answers) - 1:
                        story.append(Paragraph(f"<b>Extension Solution:</b><br/>{ans}", body_style))
                    else:
                        story.append(Paragraph(f"<b>Question {i+1}:</b><br/>{ans}", body_style))
            else:
                if str(answers).strip():
                    story.append(Paragraph(f"<b>Solutions:</b><br/>{answers}", body_style))

    doc.build(story)
    buffer.seek(0)
    return buffer.getvalue()
